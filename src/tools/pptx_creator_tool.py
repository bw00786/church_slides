
from crewai.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import json

@tool
def create_service_slides() -> str:
    """
    Generate a PowerPoint deck from structured slide data.
    This tool will be called manually with the required parameters.
    
    Returns:
        str: Confirmation message with saved PowerPoint path.
    """
    return "This tool should be called with specific parameters. Please use the manual execution method below."

def create_powerpoint_manual(slides_data, output_path, theme_backgrounds_path=None):
    """
    Manual function to create PowerPoint - call this directly from your code
    Now with automatic countdown video embedding!
    """
    print("🚀 Manual PowerPoint Creation Started")
    print(f"📝 Output path: {output_path}")
    print(f"🎨 Backgrounds path: {theme_backgrounds_path}")
    
    # Handle input data
    if isinstance(slides_data, str):
        try:
            slides_data = json.loads(slides_data)
        except:
            print("❌ Could not parse slides_data as JSON")
            return f"Error: Invalid slides_data format"
    
    if not isinstance(slides_data, list):
        print(f"❌ slides_data must be a list, got {type(slides_data)}")
        return f"Error: slides_data must be a list"
    
    print(f"📊 Processing {len(slides_data)} slides...")
    
    # Create presentation
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    successful_slides = 0
    
    for i, slide_info in enumerate(slides_data):
        if not isinstance(slide_info, dict):
            print(f"⚠️ Slide {i+1} is not a dictionary, skipping")
            continue
            
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide_info.get("title", f"Slide {i+1}")
            bg_path = slide_info.get("background_path", "")
            slide_type = slide_info.get("type", "")
            
            print(f"🖼️ Slide {i+1}: '{title}' - Background path: {bg_path}")
            
            # --- Enhanced Background image handling ---
            background_used = False
            
            if bg_path:
                possible_paths = []
                possible_paths.append(bg_path)
                
                if "/" not in bg_path and "\\" not in bg_path:
                    if theme_backgrounds_path:
                        possible_paths.append(os.path.join(theme_backgrounds_path, bg_path))
                    possible_paths.append(os.path.join("backgrounds", "forgiveness", bg_path))
                    possible_paths.append(os.path.join(os.getcwd(), "backgrounds", "forgiveness", bg_path))
                
                possible_paths.append(bg_path.replace("/", "\\"))
                possible_paths.append(bg_path.replace("\\", "/"))
                
                if theme_backgrounds_path and not os.path.isabs(bg_path):
                    possible_paths.append(os.path.join(theme_backgrounds_path, bg_path))
                    possible_paths.append(os.path.join(os.getcwd(), theme_backgrounds_path, bg_path))
                
                for test_path in set(possible_paths):
                    test_path = os.path.normpath(test_path)
                    
                    if os.path.exists(test_path):
                        try:
                            slide.shapes.add_picture(test_path, 0, 0, width=slide_width, height=slide_height)
                            print(f"✅ Slide {i+1}: '{title}' - Background FOUND: {test_path}")
                            background_used = True
                            break
                        except Exception as e:
                            print(f"⚠️ Slide {i+1}: Error loading background {test_path}: {e}")
                            continue
                
                if not background_used:
                    print(f"❌ Slide {i+1}: '{title}' - No background found after trying {len(set(possible_paths))} paths")
            
            # Fallback background
            if not background_used:
                print(f"🎨 Slide {i+1}: Using fallback blue background")
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 32, 96)
            
            # --- AUTO-EMBED COUNTDOWN VIDEO ON FIRST SLIDE ---
            if i == 0 and slide_type == 'countdown':
                # Look for countdown video
                countdown_video_path = slide_info.get('countdown_video')
                
                if not countdown_video_path:
                    # Check default location
                    countdown_video_path = 'output/countdown.mp4'
                
                if countdown_video_path and os.path.exists(countdown_video_path):
                    try:
                        # Embed video in the center of the slide
                        video_width = Inches(10)
                        video_height = Inches(5.625)  # 16:9 aspect ratio
                        video_left = (slide_width - video_width) / 2
                        video_top = (slide_height - video_height) / 2
                        
                        # Add video to slide
                        movie = slide.shapes.add_movie(
                            countdown_video_path,
                            video_left, video_top,
                            video_width, video_height
                        )
                        
                        print(f"🎬 Slide {i+1}: Countdown video EMBEDDED from {countdown_video_path}")
                        print(f"   ⚠️ Note: Video will need to be set to auto-play in PowerPoint")
                        
                    except Exception as e:
                        print(f"⚠️ Slide {i+1}: Could not embed video: {e}")
                        print(f"   You can manually insert: {countdown_video_path}")
                else:
                    print(f"⚠️ Slide {i+1}: Countdown video not found at {countdown_video_path}")
            
            # --- Text content area ---
            content_width = Inches(9)
            content_height = Inches(5)
            content_left = (slide_width - content_width) / 2
            content_top = (slide_height - content_height) / 2
            
            # Only add text content if there's content or it's not a countdown slide
            content = slide_info.get("content", "")
            if content or slide_type != 'countdown':
                # Translucent background for text
                rect = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    content_left, content_top, content_width, content_height
                )
                fill = rect.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 0, 0)
                fill.transparency = 0.3
                rect.line.fill.background()
                
                # Text box
                text_margin = Inches(0.3)
                textbox = slide.shapes.add_textbox(
                    content_left + text_margin,
                    content_top + text_margin,
                    content_width - (text_margin * 2),
                    content_height - (text_margin * 2)
                )
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                
                # Title
                p_title = text_frame.add_paragraph()
                p_title.text = str(title)
                p_title.font.bold = True
                p_title.font.size = Pt(36)
                p_title.font.color.rgb = RGBColor(255, 255, 255)
                p_title.alignment = PP_ALIGN.CENTER
                
                # Content
                if content:
                    # Add spacing after title
                    p_spacing = text_frame.add_paragraph()
                    p_spacing.text = ""
                    p_spacing.font.size = Pt(8)
                    
                    # Process content lines
                    content_lines = str(content).split("\n")
                    for line in content_lines:
                        if line.strip():
                            p = text_frame.add_paragraph()
                            p.text = line.strip()
                            p.font.size = Pt(24)
                            p.font.color.rgb = RGBColor(240, 240, 240)
                            p.alignment = PP_ALIGN.CENTER
            
            successful_slides += 1
            print(f"✅ Successfully created slide {i+1}: {title}")
            
        except Exception as e:
            print(f"❌ Error creating slide {i+1}: {e}")
            import traceback
            traceback.print_exc()
            continue
    
    # Save the presentation
    try:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        print(f"💾 Successfully saved PowerPoint with {successful_slides} slides to: {output_path}")
        return f"✅ PowerPoint created successfully with {successful_slides} slides: {output_path}"
    except Exception as e:
        error_msg = f"❌ Error saving PowerPoint: {e}"
        print(error_msg)
        return error_msg

def execute_powerpoint_creation(design_task_output, output_path, theme_backgrounds_path):
    """
    Execute PowerPoint creation directly using the design task output
    """
    print("🎯 Direct PowerPoint Execution")
    print(f"🎨 Using backgrounds from: {theme_backgrounds_path}")
    
    slides_data = None
    
    if hasattr(design_task_output, 'raw'):
        raw_output = design_task_output.raw
        if isinstance(raw_output, str):
            try:
                slides_data = json.loads(raw_output)
            except:
                import re
                json_match = re.search(r'\[.*\]', raw_output, re.DOTALL)
                if json_match:
                    try:
                        slides_data = json.loads(json_match.group())
                    except:
                        pass
        elif isinstance(raw_output, list):
            slides_data = raw_output
    
    if not slides_data:
        print("❌ Could not extract slides data from design task output")
        return False
    
    print(f"📊 Found {len(slides_data)} slides in design output")
    
    for slide in slides_data:
        if 'background_path' in slide:
            bg_path = slide['background_path']
            if bg_path and '/' not in bg_path and '\\' not in bg_path:
                new_path = os.path.join(theme_backgrounds_path, bg_path)
                print(f"🔄 Fixed background path: '{bg_path}' -> '{new_path}'")
                slide['background_path'] = new_path
    
    return create_powerpoint_manual(slides_data, output_path, theme_backgrounds_path)
